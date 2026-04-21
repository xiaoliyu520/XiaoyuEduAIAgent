from io import BytesIO
from typing import Optional

import numpy as np
from PIL import Image
from pptx import Presentation

from app.services.document_loaders.ocr import get_ocr


class PPTLoader:
    SUPPORTED_EXTENSIONS = {".ppt", ".pptx"}

    @classmethod
    def _extract_shape_text(cls, shape, resp: str, ocr) -> str:
        if shape.has_text_frame:
            resp += shape.text.strip() + "\n"
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        resp += paragraph.text.strip() + "\n"
        if shape.shape_type == 13:
            try:
                image = Image.open(BytesIO(shape.image.blob))
                if image.mode != "RGB":
                    image = image.convert("RGB")
                img_array = np.array(image)
                result, _ = ocr(img_array)
                if result:
                    ocr_result = [line[1] for line in result]
                    resp += "\n".join(ocr_result)
            except Exception:
                pass
        elif shape.shape_type == 6:
            for child_shape in shape.shapes:
                resp = cls._extract_shape_text(child_shape, resp, ocr)
        return resp

    @classmethod
    def load_from_bytes(cls, file_data: bytes) -> str:
        prs = Presentation(BytesIO(file_data))
        resp = ""
        ocr = get_ocr()

        for slide in prs.slides:
            sorted_shapes = sorted(slide.shapes, key=lambda x: (x.top, x.left))
            for shape in sorted_shapes:
                resp = cls._extract_shape_text(shape, resp, ocr)

        return resp

    @classmethod
    def load_from_file(cls, file_path: str) -> str:
        with open(file_path, "rb") as f:
            return cls.load_from_bytes(f.read())

    @classmethod
    def is_supported(cls, filename: str) -> bool:
        ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        return ext in cls.SUPPORTED_EXTENSIONS
